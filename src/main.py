import streamlit as st
import pandas as pd
import PyPDF2
import pptx
import docx
import os
import openai
import requests
import json
import logging
from io import StringIO, BytesIO
from openai import OpenAI
from reportlab.lib.pagesizes import LETTER, landscape
from reportlab.lib import colors
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.platypus import (
    SimpleDocTemplate,
    Paragraph,
    Spacer,
    Image,
    Table,
    TableStyle,
    PageBreak,
)
from reportlab.lib.units import inch
import matplotlib.pyplot as plt

# Set up logging
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

# Set your API keys here or use environment variables for security
# OpenAI API key

OPENAI_API_KEY = st.secrets["OPENAI_API_KEY"]
PERPLEXITY_API_KEY = st.secrets["PERPLEXITY_API_KEY"]

# Print API keys (Note: In production, avoid printing sensitive information)
# logger.info(f"OpenAI API Key: {OPENAI_API_KEY}")
# logger.info(f"Perplexity API Key: {PERPLEXITY_API_KEY}")

# Instantiate the OpenAI client
openai_client = OpenAI(api_key=OPENAI_API_KEY)


class PerplexityClient:
    """
    A simple client to interact with Perplexity AI's API.
    Replace the `api_url` with the actual endpoint provided by Perplexity AI.
    """

    def __init__(self, api_key):
        self.api_key = api_key
        self.api_url = "https://api.perplexity.ai/v1/chat/completions"  # Placeholder URL

    def chat_completions_create(self, model, messages, temperature=0.5, max_tokens=1000):
        headers = {
            'Content-Type': 'application/json',
            'Authorization': f'Bearer {self.api_key}',
        }

        payload = {
            'model': model,
            'messages': messages,
            'temperature': temperature,
            'max_tokens': max_tokens
        }

        try:
            response = requests.post(
                self.api_url, headers=headers, json=payload)
            response.raise_for_status()
            return response.json()
        except requests.exceptions.HTTPError as http_err:
            logger.error(f"HTTP error occurred: {http_err}")
            raise
        except Exception as err:
            logger.error(f"Other error occurred: {err}")
            raise


# Ensure all necessary items are exported
__all__ = ['openai_client', 'PerplexityClient', 'perplexity_search']


def handle_file_upload():
    """
    Handles the uploading of business documents and extracts their content.

    Returns:
        list: A list of dictionaries containing filenames and their extracted content.
    """
    uploaded_files = st.file_uploader(
        "Upload your business documents",
        type=["pdf", "pptx", "docx"],
        accept_multiple_files=True
    )
    documents = []
    if uploaded_files:
        for uploaded_file in uploaded_files:
            try:
                if uploaded_file.type == "application/pdf":
                    reader = PyPDF2.PdfReader(uploaded_file)
                    text = ''
                    for page in reader.pages:
                        extracted_text = page.extract_text()
                        if extracted_text:
                            text += extracted_text + '\n'
                elif uploaded_file.type == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
                    prs = pptx.Presentation(uploaded_file)
                    text = ''
                    for slide in prs.slides:
                        for shape in slide.shapes:
                            if hasattr(shape, "text"):
                                text += shape.text + ' '
                elif uploaded_file.type == "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
                    doc = docx.Document(uploaded_file)
                    text = '\n'.join([para.text for para in doc.paragraphs])
                else:
                    st.warning(f"Unsupported file type: {uploaded_file.type}")
                    continue
                documents.append({
                    "filename": uploaded_file.name,
                    "content": text
                })
            except Exception as e:
                st.error(f"Error processing file {uploaded_file.name}: {e}")
    return documents


def extract_key_metrics(documents):
    """
    Uses OpenAI's API to extract Market Size, Target Audience, and Competitors from the documents.

    Args:
        documents (list): List of dictionaries containing filenames and content.

    Returns:
        dict: Extracted key metrics.
    """
    combined_text = "\n\n".join([doc['content'] for doc in documents])
    prompt = (
        "Extract the following information from the business documents:\n"
        "1. Market Size\n"
        "2. Target Audience\n"
        "3. Competitors\n\n"
        "Provide the information in JSON format as follows:\n"
        "{\n"
        '  "market_size": "",\n'
        '  "target_audience": "",\n'
        '  "competitors": []\n'
        "}"
    )
    try:
        response = openai_client.chat.completions.create(
            model="gpt-4",
            messages=[
                {"role": "system", "content": "You are a business analyst."},
                {"role": "user", "content": prompt + "\n\n" + combined_text}
            ],
            temperature=0.3,
            max_tokens=500
        )
        metrics = response.choices[0].message.content
        metrics_json = json.loads(metrics)
        return metrics_json
    except Exception as e:
        st.error(f"Error extracting key metrics: {e}")
        return {}


def generate_questions(context, market_type, user_inputs):
    """
    Generates a list of questions based on the context, market type, and user inputs.

    Args:
        context (str): The context or topic for the search.
        market_type (str): The type of market analysis.
        user_inputs (dict, optional): Additional user inputs for the search.

    Returns:
        list: A list of questions to query Perplexity.
    """
    questions = [
        f"What is the current state of the {market_type} market?",
        f"What are the latest trends in the {market_type} industry?",
        f"Who are the major competitors in the {market_type} market?",
        f"What is the projected growth for the {market_type} sector over the next 5 years?"
    ]
    # Add more dynamic question generation based on user inputs if necessary
    return questions


def perplexity_search(context, market_type, user_inputs=None):
    """
    Performs a search using Perplexity AI to fetch market data.

    Args:
        context (str): The context or topic for the search.
        market_type (str): The type of market analysis.
        user_inputs (dict, optional): Additional user inputs for the search.

    Returns:
        str: The concatenated search results.
    """
    logger.info(f"Starting search for {market_type}")
    results = []
    questions = generate_questions(context, market_type, user_inputs)

    for i, question in enumerate(questions, 1):
        logger.info(f"Processing question {i} for {market_type}: {question}")
        try:
            response = perplexity_client.chat_completions_create(
                model="llama-3.1-sonar-large-128k-online",
                messages=[
                    {"role": "user", "content": question}
                ],
                temperature=0.5,
                max_tokens=1000
            )
            answer = response['choices'][0]['message']['content'].strip()

            # Clean up the answer if it's one letter per line
            if all(len(line.strip()) <= 1 for line in answer.split('\n')):
                answer = ''.join(line.strip()
                                 for line in answer.split('\n') if line.strip())

            logger.info(f"Received answer: {answer[:100]}...")
            results.append(f"**Question:** {question}\n**Answer:** {answer}")
        except Exception as e:
            logger.error(f"Error fetching search results: {e}")
            results.append(
                f"**Question:** {question}\n**Answer:** Error fetching data.")

    logger.info(f"Completed search for {market_type}")
    return "\n\n".join(results)


def prepare_prompt(template, user_inputs):
    prompt = template
    for placeholder, value in user_inputs.items():
        prompt = prompt.replace(f"[{placeholder}]", value)
    return prompt


def set_research_context(document_text, user_inputs):
    initial_analysis_request = "Analyze the following business documents and provide insights."
    prompt = prepare_prompt(initial_analysis_request, user_inputs)
    prompt += f"\n\nDocument Text:\n{document_text}"

    # Your existing code to send the prompt to the AI model
    # ...


@st.cache_resource
def initialize_perplexity_client():
    """
    Initializes the Perplexity client.

    Returns:
        PerplexityClient: An instance of the Perplexity client.
    """
    # Initialize your Perplexity client here
    return PerplexityClient(api_key=PERPLEXITY_API_KEY)


# Initialize Perplexity client
perplexity_client = initialize_perplexity_client()


@st.cache_data
def perform_market_analysis(document_text, market_type, user_inputs=None):
    """
    Performs market analysis using Perplexity AI based on the document content.
    """
    logger.info(f"Starting market analysis for {market_type}")
    try:
        results = perplexity_search(document_text, market_type, user_inputs)
        logger.info(f"Completed market analysis for {market_type}")
        return results
    except Exception as e:
        logger.error(
            f"Error in perform_market_analysis for {market_type}: {e}")
        st.error(
            f"An error occurred during {market_type} analysis. Please try again later.")
        return ""


def generate_executive_summary(all_analysis_results, user_inputs):
    """
    Generates an executive summary based on the analysis results.

    Args:
        all_analysis_results (str): The concatenated analysis results.
        user_inputs (dict): Additional user inputs for context.

    Returns:
        str: The executive summary.
    """
    logger.info("Generating executive summary")
    try:
        prompt = f"""
        Based on the following comprehensive analysis results, create a detailed executive summary:

        {all_analysis_results}

        Consider the following user inputs for additional context:
        """
        for key, value in user_inputs.items():
            prompt += f"{key}: {value}\n"

        prompt += """
        Synthesize all the information provided to create a comprehensive executive summary that includes:
        1. An overview of the AI-powered business development tools market, focusing on the specific geography and target audience analyzed.
        2. Key findings from the TAM, SAM, and SOM analysis, including any projections made.
        3. Critical market trends, growth drivers, and potential challenges identified in the analysis.
        4. Competitive landscape insights and the company's potential positioning.
        5. Strategic recommendations based on the analysis.
        6. Any other significant insights or conclusions drawn from the comprehensive analysis.

        Ensure the summary is coherent, data-driven, and directly reflects the analysis performed. Do not include any information that wasn't part of the provided analysis results.
        """

        response = openai_client.chat.completions.create(
            model="gpt-4",
            messages=[
                {"role": "system", "content": "You are an AI assistant specialized in synthesizing complex market analysis into coherent executive summaries. Focus on the data and insights provided, without adding external information."},
                {"role": "user", "content": prompt}
            ],
            max_tokens=4000,
            temperature=0.7
        )
        exec_summary = response.choices[0].message.content.strip()
    except Exception as e:
        logger.error(f"Error in generate_executive_summary: {e}")
        st.error(
            "An error occurred while generating the executive summary. Please try again later.")
        return ""
    logger.info("Executive summary generated")
    return exec_summary


def extract_industry_keywords(metrics):
    """
    Derives industry-specific keywords from the extracted metrics.

    Args:
        metrics (dict): Extracted key metrics.

    Returns:
        list: List of industry-specific keywords.
    """
    keywords = []
    if 'market_size' in metrics and metrics['market_size']:
        keywords.append("market size")
    if 'target_audience' in metrics and metrics['target_audience']:
        keywords.append("target audience")
    if 'competitors' in metrics and metrics['competitors']:
        keywords.extend(metrics['competitors'])
    # Add more sophisticated keyword extraction if needed
    return keywords


def calculate_TAM_SAM_SOM(metrics, external_data):
    """
    Calculates TAM, SAM, and SOM based on extracted and external data.

    Args:
        metrics (dict): Extracted key metrics.
        external_data (dict): Data fetched from Perplexity API.

    Returns:
        tuple: TAM, SAM, SOM values.
    """
    try:
        # Calculate TAM
        if 'market_size' in metrics and metrics['market_size']:
            tam = parse_market_size(metrics['market_size'])
        else:
            # Fetch market size from external data if not available
            tam = 0
            if 'market size' in external_data:
                tam = parse_market_size(external_data['market size'])
            else:
                tam = 0  # Default to 0 if not available

        # Calculate SAM
        if 'target_audience' in metrics and metrics['target_audience']:
            target_audience_percentage = parse_target_audience(
                metrics['target_audience'])
        else:
            # Fetch target audience data from external data if not available
            if 'target audience' in external_data:
                target_audience_percentage = parse_target_audience(
                    external_data['target audience'])
            else:
                target_audience_percentage = 50  # Default value if not available

        sam = tam * (target_audience_percentage / 100)

        # Calculate SOM
        if 'market_share_estimate' in metrics and metrics['market_share_estimate']:
            market_share_percentage = parse_market_share(
                metrics['market_share_estimate'])
        else:
            # Fetch market share estimate from external data if not available
            if 'market share estimate' in external_data:
                market_share_percentage = parse_market_share(
                    external_data['market share estimate'])
            else:
                market_share_percentage = 5  # Default value if not available

        som = sam * (market_share_percentage / 100)

        return tam, sam, som
    except Exception as e:
        st.error(f"Error calculating TAM/SAM/SOM: {e}")
        return 0, 0, 0


def parse_market_size(market_size_text):
    """
    Parses the market size from text to extract numerical value.

    Args:
        market_size_text (str): Text containing market size information.

    Returns:
        float: Market size in USD.
    """
    # Implement actual parsing logic using regex or NLP
    # Enhanced implementation with more robust patterns
    try:
        import re
        patterns = [
            r'([\d,.]+)\s*(million|billion|trillion)\s*USD',
            r'([\d,.]+)\s*(million|billion|trillion)'
        ]
        for pattern in patterns:
            match = re.search(pattern, market_size_text, re.IGNORECASE)
            if match:
                number = float(match.group(1).replace(',', ''))
                scale = match.group(2).lower()
                if scale == 'million':
                    return number * 1_000_000
                elif scale == 'billion':
                    return number * 1_000_000_000
                elif scale == 'trillion':
                    return number * 1_000_000_000_000
        return 0
    except:
        return 0


def parse_target_audience(audience_text):
    """
    Parses target audience information to extract percentage.

    Args:
        audience_text (str): Text containing target audience information.

    Returns:
        float: Percentage of the market.
    """
    # Implement actual parsing logic
    # Enhanced implementation with more robust patterns
    try:
        import re
        match = re.search(r'(\d+\.?\d*)\s*%', audience_text)
        if match:
            percentage = float(match.group(1))
            return percentage if 0 <= percentage <= 100 else 50  # Validate percentage
        return 50  # Default to 50% if not specified
    except:
        return 50


def parse_market_share(market_share_text):
    """
    Parses market share information to extract percentage.

    Args:
        market_share_text (str): Text containing market share information.

    Returns:
        float: Market share percentage.
    """
    # Implement actual parsing logic
    try:
        import re
        match = re.search(r'(\d+\.?\d*)\s*%', market_share_text)
        if match:
            percentage = float(match.group(1))
            return percentage if 0 <= percentage <= 100 else 5  # Validate percentage
        return 5  # Default to 5% if not specified
    except:
        return 5


def generate_visualizations(metrics):
    """
    Generates visualizations based on the metrics.

    Args:
        metrics (dict): Extracted key metrics.

    Returns:
        dict: Dictionary containing paths to the generated images.
    """
    images = {}

    # Example: Bar chart for Market Size, SAM, SOM
    try:
        categories = ['TAM', 'SAM', 'SOM']
        values = [
            metrics.get('tam', 0),
            metrics.get('sam', 0),
            metrics.get('som', 0)
        ]

        plt.figure(figsize=(6, 4))
        plt.bar(categories, values, color=['blue', 'green', 'red'])
        plt.title('Market Size Analysis')
        plt.xlabel('Categories')
        plt.ylabel('Amount (USD)')
        plt.tight_layout()

        buf = BytesIO()
        plt.savefig(buf, format='PNG')
        buf.seek(0)
        images['market_size_chart'] = buf
        plt.close()
    except Exception as e:
        logger.error(f"Error generating market size chart: {e}")

    # Add more visualizations as needed

    return images


def generate_pdf_report(extracted_info, tam, sam, som, executive_summary, images):
    """
    Generates a detailed PDF report with visualizations and structured sections.

    Args:
        extracted_info (dict): Extracted key metrics.
        tam (float): Total Addressable Market.
        sam (float): Serviceable Available Market.
        som (float): Serviceable Obtainable Market.
        executive_summary (str): Executive summary text.
        images (dict): Dictionary of image buffers.

    Returns:
        BytesIO: In-memory PDF file.
    """
    buffer = BytesIO()
    doc = SimpleDocTemplate(buffer, pagesize=LETTER)
    story = []
    styles = getSampleStyleSheet()
    styles.add(ParagraphStyle(name='CenteredTitle',
               alignment=1, fontSize=24, spaceAfter=20))
    styles.add(ParagraphStyle(name='SectionHeading', fontSize=18,
               spaceAfter=10, textColor=colors.HexColor("#2E86AB")))
    styles.add(ParagraphStyle(name='NormalText', fontSize=12, spaceAfter=10))

    # Cover Page
    story.append(Spacer(1, 2 * inch))
    story.append(Paragraph("Business Document Analysis Report",
                 styles['CenteredTitle']))
    story.append(Spacer(1, 0.5 * inch))
    story.append(Paragraph("Generated by AI-Powered Analyzer",
                 styles['CenteredTitle']))
    story.append(PageBreak())

    # Introduction
    story.append(Paragraph("Introduction", styles['SectionHeading']))
    intro_text = (
        "This report provides a comprehensive analysis of the uploaded business documents, "
        "extracting key metrics such as Market Size, Target Audience, and Competitors. "
        "Further market analysis was conducted using Perplexity AI to derive additional insights "
        "and forecasts."
    )
    story.append(Paragraph(intro_text, styles['NormalText']))
    story.append(Spacer(1, 0.2 * inch))

    # Extracted Key Metrics
    story.append(Paragraph("Extracted Key Metrics", styles['SectionHeading']))
    data = [
        ['Market Size', extracted_info.get('market_size', 'N/A')],
        ['Target Audience', extracted_info.get('target_audience', 'N/A')],
        ['Competitors', ', '.join(extracted_info.get('competitors', []))]
    ]
    table = Table(data, colWidths=[2 * inch, 4 * inch])
    table.setStyle(TableStyle([
        ('BACKGROUND', (0, 0), (-1, 0), colors.HexColor("#D6EAF8")),
        ('TEXTCOLOR', (0, 0), (-1, 0), colors.black),
        ('ALIGN', (0, 0), (-1, -1), 'LEFT'),
        ('GRID', (0, 0), (-1, -1), 1, colors.black),
        ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
    ]))
    story.append(table)
    story.append(Spacer(1, 0.3 * inch))

    # Market Analysis
    story.append(Paragraph("Market Analysis", styles['SectionHeading']))
    analysis_text = (
        f"- **TAM (Total Addressable Market):** ${tam:,.2f}\n"
        f"- **SAM (Serviceable Available Market):** ${sam:,.2f}\n"
        f"- **SOM (Serviceable Obtainable Market):** ${som:,.2f}"
    )
    story.append(Paragraph(analysis_text, styles['NormalText']))
    story.append(Spacer(1, 0.2 * inch))

    # Visualization: Market Size Analysis
    if 'market_size_chart' in images:
        story.append(Paragraph("Market Size Analysis Chart",
                     styles['SectionHeading']))
        img = Image(images['market_size_chart'], width=6*inch, height=4*inch)
        story.append(img)
        story.append(Spacer(1, 0.3 * inch))

    # Executive Summary
    story.append(PageBreak())
    story.append(Paragraph("Executive Summary", styles['SectionHeading']))
    exec_summary_paragraphs = executive_summary.split('\n')
    for para in exec_summary_paragraphs:
        story.append(Paragraph(para, styles['NormalText']))
    story.append(Spacer(1, 0.2 * inch))

    # Add more sections as needed (e.g., Insights, Recommendations)

    # Build PDF
    doc.build(story)
    buffer.seek(0)
    return buffer


def generate_report(extracted_info, tam, sam, som):
    """
    Generates a report based on the extracted information and market calculations.

    Args:
        extracted_info (dict): Extracted key metrics.
        tam (float): Total Addressable Market.
        sam (float): Serviceable Available Market.
        som (float): Serviceable Obtainable Market.

    Returns:
        str: Formatted report as a string.
    """
    try:
        report = "## Business Document Analysis Report\n\n"
        report += "### Extracted Key Metrics\n"
        report += f"- **Market Size:** {extracted_info.get('market_size', 'N/A')}\n"
        report += f"- **Target Audience:** {extracted_info.get('target_audience', 'N/A')}\n"
        report += f"- **Market Share Estimate:** {extracted_info.get('market_share_estimate', 'N/A')}\n"
        report += f"- **Competitors:** {', '.join(extracted_info.get('competitors', []))}\n\n"

        report += "### Market Analysis\n"
        report += f"- **TAM (Total Addressable Market):** ${tam:,.2f}\n"
        report += f"- **SAM (Serviceable Available Market):** ${sam:,.2f}\n"
        report += f"- **SOM (Serviceable Obtainable Market):** ${som:,.2f}\n\n"

        report += "### Insights\n"
        report += "Based on the extracted data and external market research, the above metrics provide a comprehensive overview of the market potential and strategic positioning.\n"

        return report
    except Exception as e:
        st.error(f"Error generating report: {e}")
        return "Error generating report."


def fetch_additional_metrics(metrics, external_data):
    """
    Ensures all necessary metrics are available by fetching missing data from external sources.

    Args:
        metrics (dict): Extracted key metrics.
        external_data (dict): Data fetched from Perplexity API.

    Returns:
        dict: Complete metrics with any missing data filled from external sources.
    """
    complete_metrics = metrics.copy()

    # Ensure market_size
    if not complete_metrics.get('market_size'):
        if 'market size' in external_data:
            complete_metrics['market_size'] = external_data['market size']
        else:
            st.warning(
                "Market size data is missing and could not be fetched externally.")

    # Ensure target_audience
    if not complete_metrics.get('target_audience'):
        if 'target audience' in external_data:
            complete_metrics['target_audience'] = external_data['target audience']
        else:
            st.warning(
                "Target audience data is missing and could not be fetched externally.")

    # Ensure market_share_estimate
    if not complete_metrics.get('market_share_estimate'):
        if 'market share estimate' in external_data:
            complete_metrics['market_share_estimate'] = external_data['market share estimate']
        else:
            complete_metrics['market_share_estimate'] = "5%"  # Default value

    # Ensure competitors
    if not complete_metrics.get('competitors'):
        if 'competitors' in external_data and isinstance(external_data['competitors'], list):
            complete_metrics['competitors'] = external_data['competitors']
        else:
            st.warning(
                "Competitors data is missing and could not be fetched externally.")
            complete_metrics['competitors'] = []

    return complete_metrics


def parse_external_data(analysis_result):
    """
    Parses the external analysis result from Perplexity into a dictionary.

    Args:
        analysis_result (str): The raw analysis result from Perplexity.

    Returns:
        dict: Parsed external data.
    """
    external_data = {}
    try:
        # Example parsing logic; adjust based on actual Perplexity response format
        lines = analysis_result.split('\n')
        for line in lines:
            if "Market Size" in line:
                external_data['market size'] = line.split(
                    "Answer:")[-1].strip()
            elif "Target Audience" in line:
                external_data['target audience'] = line.split(
                    "Answer:")[-1].strip()
            elif "Market Share Estimate" in line:
                external_data['market share estimate'] = line.split(
                    "Answer:")[-1].strip()
            elif "Competitors" in line:
                competitors = line.split("Answer:")[-1].strip()
                external_data['competitors'] = [c.strip()
                                                for c in competitors.split(',') if c.strip()]
    except Exception as e:
        logger.error(f"Error parsing external data: {e}")
    return external_data


def main():
    """
    Main function to run the Streamlit app.
    """
    st.title("📄 Business Documentation Analyzer")
    st.write(
        "Upload your business documents to analyze and generate a comprehensive PDF report.")

    # File Upload
    documents = handle_file_upload()

    if documents:
        st.success("Files successfully uploaded and processed.")

        # Display Uploaded Files
        st.subheader("Uploaded Documents")
        for doc in documents:
            st.write(f"- {doc['filename']}")

        # Extract Key Metrics using OpenAI
        with st.spinner("Extracting key metrics using OpenAI..."):
            key_metrics = extract_key_metrics(documents)

        if key_metrics:
            st.subheader("Extracted Key Metrics")
            st.json(key_metrics)

            # Derive Industry-Specific Keywords
            industry_keywords = extract_industry_keywords(key_metrics)
            st.write("### Industry-Specific Keywords")
            st.write(", ".join(industry_keywords))

            # Query Perplexity API for Additional Market Data
            external_data = {}
            with st.spinner("Fetching additional market data from Perplexity..."):
                for keyword in industry_keywords:
                    market_type = keyword  # Assuming each keyword represents a market type
                    user_inputs = {"keyword": keyword}
                    analysis_result = perform_market_analysis(
                        keyword, market_type, user_inputs)
                    if analysis_result:
                        external_data.update(
                            parse_external_data(analysis_result))

            # Ensure all metrics are complete
            key_metrics = fetch_additional_metrics(key_metrics, external_data)

            st.subheader("Complete Key Metrics")
            st.json(key_metrics)

            st.subheader("External Market Data")
            st.json(external_data)

            # Calculate TAM, SAM, SOM
            tam, sam, som = calculate_TAM_SAM_SOM(key_metrics, external_data)
            st.subheader("Market Analysis")
            st.write(f"**TAM (Total Addressable Market):** ${tam:,.2f}")
            st.write(f"**SAM (Serviceable Available Market):** ${sam:,.2f}")
            st.write(f"**SOM (Serviceable Obtainable Market):** ${som:,.2f}")

            # Generate Report Text (for display and PDF generation)
            report_text = generate_report(key_metrics, tam, sam, som)
            st.subheader("Analysis Report")
            st.markdown(report_text)

            # Generate Executive Summary
            executive_summary = generate_executive_summary(
                report_text, key_metrics)
            if executive_summary:
               #  st.subheader("Executive Summary")
                st.markdown(executive_summary)

            # Generate Visualizations
            images = generate_visualizations({
                'tam': tam,
                'sam': sam,
                'som': som,
                'market_size': key_metrics.get('market_size'),
                'target_audience': key_metrics.get('target_audience'),
                'competitors': key_metrics.get('competitors'),
            })

            # Generate PDF
            with st.spinner("Generating PDF report..."):
                pdf_buffer = generate_pdf_report(
                    extracted_info=key_metrics,
                    tam=tam,
                    sam=sam,
                    som=som,
                    executive_summary=executive_summary,
                    images=images
                )

            # Provide Download Button for PDF
            st.download_button(
                label="📥 Download PDF Report",
                data=pdf_buffer,
                file_name="analysis_report.pdf",
                mime="application/pdf"
            )
        else:
            st.error("Failed to extract key metrics from the documents.")
    else:
        st.info("Awaiting document uploads.")


if __name__ == "__main__":
    main()
